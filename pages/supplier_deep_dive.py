import streamlit as st
import pandas as pd
import numpy as np
import plotly.graph_objects as go
import plotly.express as px
import plotly.figure_factory as ff
from sklearn.ensemble import RandomForestClassifier
from prophet import Prophet
from pptx import Presentation
from pptx.util import Inches
import io
import warnings

# --- SUPPRESS DEPRECATION WARNINGS FOR CLEANER PRESENTATION ---
warnings.filterwarnings("ignore", category=DeprecationWarning)
warnings.filterwarnings("ignore", category=FutureWarning)


# --- ROBUST STATE CHECK ---
if 'app_data' not in st.session_state:
    st.error("Application data not loaded. Please return to the 'Global Command Center' home page to initialize the app.")
    st.stop()

# Unpack data
app_data = st.session_state['app_data']
suppliers = app_data['suppliers']
foundry_perf = app_data['foundry_perf']
osat_perf = app_data['osat_perf']
failures = app_data['failures']

# --- UI RENDER ---
st.markdown("# 🔬 Supplier Deep Dive & Process Control")
st.markdown("This module is the SQE's workbench for monitoring historical performance, analyzing process capability, predicting future outcomes, and taking formal corrective action. The tools shown are **dynamically adapted** based on the selected supplier's type (Foundry or OSAT).")

selected_supplier = st.selectbox("Select a Supplier to Analyze", suppliers['Supplier'].unique(), key="supplier_select_deep_dive")
supplier_info = suppliers[suppliers['Supplier'] == selected_supplier].iloc[0]
supplier_type = supplier_info['Type']

st.info(f"**Viewing:** `{selected_supplier}` | **Supplier Type:** `{supplier_type}`. The analytical tools below are tailored for this supplier type.", icon="💡")

tab_monitor, tab_predict, tab_act = st.tabs(["📊 Process Monitoring (SPC & Cpk)", "🔮 Predictive Analytics (Forecast & ML)", "✍️ Corrective Action & Reporting (SCAR)"])

with tab_monitor:
    st.header(f"Process Monitoring for: {selected_supplier} ({supplier_type})")

    if supplier_type == 'Foundry':
        st.markdown("Analyzing critical **Frontend (Wafer Fab)** process control and capability metrics.")
        col1, col2 = st.columns(2)
        with col1:
            st.subheader("Wafer Acceptance Test (WAT) Stability")
            st.markdown("- **What:** An SPC chart on a critical transistor parameter, `Threshold Voltage (Vt)`, measured on test structures on every wafer. \n- **Why:** WAT is the gatekeeper for foundry quality. A stable Vt is essential for the ASIC's performance and power consumption. An out-of-control point here is a leading indicator of a major process excursion. \n- **Standard:** **JEDEC JESD47** (Stress-Test Driven Qualification).")
            np.random.seed(hash(selected_supplier) % (2**32 - 1)); wat_data = np.random.normal(loc=0.45, scale=0.01, size=50)
            fig_spc = go.Figure(); fig_spc.add_trace(go.Scatter(y=wat_data, mode='lines+markers', name='Vt Measurement'))
            fig_spc.add_hline(y=0.45, line=dict(dash="dash", color="green"), name="Target"); fig_spc.add_hline(y=0.48, line=dict(dash="dot", color="red"), name="UCL"); fig_spc.add_hline(y=0.42, line=dict(dash="dot", color="red"), name="LCL")
            fig_spc.update_layout(title="SPC on Threshold Voltage (Vt)", yaxis_title="Voltage (V)", xaxis_title="Wafer Lot")
            st.plotly_chart(fig_spc, use_container_width=True, key="foundry_spc_chart")
        with col2:
            st.subheader("Process Capability (Cpk) for Vt")
            st.markdown("- **Why (Actionability):** For ASICs, process capability is paramount. A Cpk below 1.33 means the foundry process is not robust enough and will produce dies that fail at different operating conditions (process corners). This is a data-driven basis for rejecting a wafer lot or demanding process improvement. \n- **Standard:** A core **Six Sigma** metric, essential for **AS9145 (APQP)** process validation.")
            usl, lsl = 0.5, 0.4; mu, sigma = 0.455, 0.015; process_data = np.random.normal(mu, sigma, 200)
            cpu = (usl - mu) / (3 * sigma); cpl = (mu - lsl) / (3 * sigma); cpk = min(cpu, cpl)
            fig_cpk = ff.create_distplot([process_data], ['Vt Data'], show_hist=True, show_rug=False)
            fig_cpk.add_vline(x=usl, line=dict(dash="dash", color="red"), name="USL"); fig_cpk.add_vline(x=lsl, line=dict(dash="dash", color="red"), name="LSL")
            fig_cpk.update_layout(title=f"Process Capability: Threshold Voltage (Cpk = {cpk:.2f})")
            st.plotly_chart(fig_cpk, use_container_width=True, key="foundry_cpk_chart")

        # --- ENHANCED VISUALIZATION: Wafer Defect Map with Scale ---
        st.subheader("Wafer Defect Bin Map")
        st.markdown("""
        - **What:** A heatmap visualizing the **Bin Code** of each die on a simulated 300mm wafer. Bin 1 (green) represents a good die, while other colored bins represent specific failure modes.
        - **How:** A 2D array is generated to simulate die locations. A circular mask is applied for realism. The color and the **legend on the right (the scale)** map each numerical bin to a specific failure category.
        - **Why (Actionability):** This is a critical diagnostic tool. Randomly scattered defects ("salt-and-pepper") are expected. However, repeating patterns, like the **red "edge ring" of Bin 4 failures** shown here, are 'smoking guns' that point to specific equipment problems in the fab (e.g., photoresist spin-coater issues or plasma etcher non-uniformity).
        """)
        
        @st.cache_data
        def generate_wafer_map():
            size = 50
            # Start with all good dies (Bin 1)
            wafer_map = np.ones((size, size), dtype=int)
            
            # Create a circular mask (set dies outside the circle to Bin 0: No Die)
            center = size / 2
            radius = size / 2 - 2
            for i in range(size):
                for j in range(size):
                    if (i - center)**2 + (j - center)**2 > radius**2:
                        wafer_map[i, j] = 0 # Bin 0 = No Die / Edge Exclusion
                    # Create a systematic "edge ring" defect pattern (Bin 4)
                    elif radius**2 > (i - center)**2 + (j - center)**2 > (radius - 3)**2:
                        if np.random.rand() > 0.6: # 40% of edge dies fail
                           wafer_map[i, j] = 4
            
            # Sprinkle some random defects
            num_random_defects = int(size * size * 0.02) # 2% random defects
            for _ in range(num_random_defects):
                x, y = np.random.randint(0, size, 2)
                if wafer_map[x, y] == 1: # Only place defects on good dies
                    wafer_map[x, y] = np.random.choice([2, 3]) # Bin 2 or 3
            return wafer_map
        
        wafer_map_data = generate_wafer_map()
        
        # Define the scale: map bin numbers to colors and descriptive labels
        color_map = {0: 'lightgrey', 1: 'mediumseagreen', 2: 'orange', 3: 'yellow', 4: 'indianred'}
        bin_labels = {0: 'No Die', 1: 'Bin 1: Good', 2: 'Bin 2: Leakage Fail', 3: 'Bin 3: Speed Fail', 4: 'Bin 4: Edge Defect'}

        # Create a DataFrame for Plotly to easily use labels
        df_wafer = pd.DataFrame(wafer_map_data).stack().reset_index()
        df_wafer.columns = ['Y', 'X', 'Bin']
        df_wafer['Label'] = df_wafer['Bin'].map(bin_labels)

        fig_wafer = px.imshow(
            wafer_map_data,
            color_continuous_scale=[color_map[i] for i in sorted(color_map.keys())],
            title="Wafer Sort Yield Map - Lot XA-123",
            labels=dict(color="Bin Code")
        )
        
        # Update the color bar (the scale) with proper labels
        fig_wafer.update_coloraxes(
            colorbar=dict(
                title="Bin Legend",
                tickvals=sorted(bin_labels.keys()),
                ticktext=list(bin_labels.values()),
            )
        )
        st.plotly_chart(fig_wafer, use_container_width=True, key="wafer_map")

    else: # OSAT Section
        # This section is preserved from the previous version
        st.markdown("Analyzing critical **Backend (Assembly & Test)** process control and capability metrics.")
        col1, col2 = st.columns(2)
        with col1:
            st.subheader("Assembly Process Stability (Wire Bond)")
            st.markdown("- **What:** An SPC chart monitoring the shear strength of wire bonds, a critical mechanical test. \n- **Why:** For a satellite that must survive launch vibrations, mechanical integrity is paramount. A drop in wire bond strength is a major reliability risk. \n- **Standard:** **MIL-STD-883 Test Method 2011**, **JEDEC JESD22-B116**.")
            np.random.seed(hash(selected_supplier) % (2**32 - 1)); shear_data = np.random.normal(loc=8.5, scale=0.2, size=50)
            fig_spc_osat = go.Figure(); fig_spc_osat.add_trace(go.Scatter(y=shear_data, mode='lines+markers', name='Shear Strength'))
            fig_spc_osat.add_hline(y=8.5, line=dict(dash="dash", color="green"), name="Target"); fig_spc_osat.add_hline(y=9.1, line=dict(dash="dot", color="red"), name="UCL"); fig_spc_osat.add_hline(y=7.9, line=dict(dash="dot", color="red"), name="LCL")
            fig_spc_osat.update_layout(title="SPC on Wire Bond Shear Strength", yaxis_title="Force (grams)", xaxis_title="Assembly Lot")
            st.plotly_chart(fig_spc_osat, use_container_width=True, key="osat_spc_chart")
        with col2:
            st.subheader("Final Test Bin-Out Pareto")
            st.markdown("- **Why (Actionability):** This is the most important chart for diagnosing test failures at an OSAT. It immediately tells the SQE where to focus. A high count in 'Continuity/Opens' points to an assembly problem, while a high count in 'Max Frequency' points to a silicon performance issue. \n- **Standard:** Data is collected per **IPC-9261** (Assembly Process Monitoring).")
            bin_data = pd.DataFrame({'Bin': ['Bin 1: Good', 'Bin 2: Continuity/Opens', 'Bin 5: Max Freq Fail', 'Bin 8: IO Leakage', 'Bin 3: Shorts'], 'Count': [9850, 75, 45, 22, 8]})
            fig_pareto_osat = px.bar(bin_data[bin_data['Bin'] != 'Bin 1: Good'], x='Count', y='Bin', orientation='h', title="Final Test Bin-Out Failures - Lot #7891", text='Count')
            fig_pareto_osat.update_layout(yaxis={'categoryorder':'total ascending'})
            st.plotly_chart(fig_pareto_osat, use_container_width=True, key="osat_pareto_chart")

# The rest of the page (Predictive Analytics and SCAR Reporting) is preserved as it is already robust and context-aware.
with tab_predict:
    # This code is preserved from the previous version...
    st.header(f"Predictive Analytics for: {selected_supplier} ({supplier_type})")
    st.subheader("Future Performance Forecast (Prophet)")
    if supplier_type == 'Foundry':
        forecast_metric = 'Wafer_Sort_Yield'
        perf_data_for_forecast = foundry_perf[foundry_perf['Supplier'] == selected_supplier].copy()
        st.markdown("- **Why:** Forecasting wafer sort yield helps predict the raw silicon supply for the entire downstream chain. A forecasted dip here will impact OSATs and final satellite production weeks later.")
    else: # OSAT
        forecast_metric = 'DPPM'
        perf_data_for_forecast = osat_perf[osat_perf['Supplier'] == selected_supplier].copy()
        st.markdown("- **Why:** Forecasting DPPM helps predict the quality of parts arriving at Kuiper. A forecasted spike allows us to proactively increase incoming inspection or put the supplier on notice.")
    @st.cache_data
    def run_prophet_forecast(data, metric_col, periods=30):
        m = Prophet(daily_seasonality=False, weekly_seasonality=True, yearly_seasonality=True); m.fit(data[['Date', metric_col]].rename(columns={'Date': 'ds', metric_col: 'y'}))
        return m, m.predict(m.make_future_dataframe(periods=periods))
    model_prophet, forecast = run_prophet_forecast(perf_data_for_forecast, forecast_metric)
    fig_forecast = go.Figure(); fig_forecast.add_trace(go.Scatter(x=forecast['ds'], y=forecast['yhat'], mode='lines', name='Forecast', line=dict(color='navy', dash='dash')))
    fig_forecast.add_trace(go.Scatter(x=forecast['ds'], y=forecast['yhat_upper'], fill=None, mode='lines', line_color='rgba(0,176,246,0.2)', name='Uncertainty')); fig_forecast.add_trace(go.Scatter(x=forecast['ds'], y=forecast['yhat_lower'], fill='tonexty', mode='lines', line_color='rgba(0,176,246,0.2)'))
    fig_forecast.add_trace(go.Scatter(x=model_prophet.history['ds'], y=model_prophet.history['y'], mode='markers', name='Actuals', marker=dict(color='black', size=4)))
    fig_forecast.update_layout(title=f"30-Day {forecast_metric.replace('_', ' ')} Forecast", yaxis_title=forecast_metric)
    st.plotly_chart(fig_forecast, use_container_width=True, key="prophet_forecast_chart_context")
    st.subheader("Predictive Lot Disposition (ML Classifier)")
    st.markdown("- **Why:** This enables a 'smarter' incoming inspection (IQC) strategy. We can allocate more stringent testing to lots the model flags as high-risk, optimizing resources and improving escape detection.")
    @st.cache_data
    def get_model_and_data():
        np.random.seed(42); X = pd.DataFrame({'Temp_Avg': np.random.normal(150, 5, 200), 'Pressure_Var': np.random.gamma(1, 0.5, 200), 'Vibration_Max': np.random.uniform(0.1, 1.0, 200)})
        y = ((X['Temp_Avg'] > 155) | (X['Pressure_Var'] > 1.2) | (X['Vibration_Max'] > 0.8)).astype(int); y = y & (np.random.rand(200) < 0.7)
        model = RandomForestClassifier(n_estimators=50, random_state=42).fit(X, y); return model, X.describe()
    model_rf, X_desc = get_model_and_data()
    col1, col2 = st.columns([1, 2])
    with col1:
        temp = st.slider("Average Temp (°C)", float(X_desc.loc['min','Temp_Avg']), float(X_desc.loc['max','Temp_Avg']), 152.0, 0.1, key="slider_temp")
        pressure = st.slider("Pressure Variance (psi)", float(X_desc.loc['min','Pressure_Var']), float(X_desc.loc['max','Pressure_Var']), 0.8, 0.01, key="slider_pressure")
        vibration = st.slider("Max Vibration (g)", float(X_desc.loc['min','Vibration_Max']), float(X_desc.loc['max','Vibration_Max']), 0.5, 0.01, key="slider_vibration")
    with col2:
        input_data = pd.DataFrame([[temp, pressure, vibration]], columns=['Temp_Avg', 'Pressure_Var', 'Vibration_Max']); fail_prob = model_rf.predict_proba(input_data)[0, 1]
        if fail_prob > 0.6: st.error(f"**High Risk ({fail_prob:.0%})** - Recommend placing lot on hold for engineering review.", icon="🚨")
        elif fail_prob > 0.3: st.warning(f"**Medium Risk ({fail_prob:.0%})** - Recommend enhanced inspection (per **ANSI Z1.4**).", icon="⚠️")
        else: st.success(f"**Low Risk ({fail_prob:.0%})** - Recommend standard release protocol.", icon="✅")
        st.progress(fail_prob)

with tab_act:
    st.header(f"Formal Corrective Action for: {selected_supplier} ({supplier_type})")
    st.markdown("- **Why:** This automates a critical communication workflow. A formal SCAR is a key part of any **ISO 9001 / AS9100** compliant quality system.")
    st.subheader("Generate Supplier Corrective Action Request (SCAR)")
    if supplier_type == 'Foundry':
        default_desc = "During Wafer Sort of Lot WN-45B, a cluster of dies in the upper-left quadrant failed parametric tests for Vt, indicating a potential process non-uniformity."
        default_standards = ["AS9100D Clause 8.5.1: Control of production and service provision", "JEDEC JESD47: Stress-Test Driven Qualification of ICs"]
    else: # OSAT
        default_desc = "During incoming inspection of Lot A-LOT-7891, 15 out of 1000 units exhibited wire bond shear failures below the 5-gram engineering limit."
        default_standards = ["AS9100D Clause 8.4.3: Information for External Providers", "IPC-A-610 Class 3: Acceptance Criteria for Wire Bonds"]
    non_conformance_details = st.text_area("Describe the Non-Conformance (Be specific)", default_desc)
    standard_ref = st.selectbox("Reference to Quality Standard Requirement", default_standards)
    if st.button("Generate SCAR PowerPoint"):
        with st.spinner("Creating SCAR..."):
            prs = Presentation(); slide = prs.slides.add_slide(prs.slide_layouts[0]); slide.shapes.title.text = "Supplier Corrective Action Request (SCAR)"; slide.placeholders[1].text = f"To: {selected_supplier}\nSCAR ID: KUI-SCAR-2023-018"
            slide = prs.slides.add_slide(prs.slide_layouts[5]); slide.shapes.title.text = "SCAR Details & Objective Evidence"
            ppt_buffer = io.BytesIO(); prs.save(ppt_buffer); ppt_buffer.seek(0)
            st.success("SCAR generated successfully!")
            st.download_button(label="Download SCAR (.pptx)", data=ppt_buffer, file_name=f"SCAR_{selected_supplier.replace(' ', '_')}_{pd.Timestamp.now().strftime('%Y%m%d')}.pptx")
